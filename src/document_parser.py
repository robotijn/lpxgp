"""Document parsing utilities for pitch deck text extraction.

This module provides functions to extract text from PDF and PowerPoint files,
primarily used for parsing pitch decks uploaded by GPs.

Supported formats:
- PDF (.pdf) - using pdfplumber
- PowerPoint (.pptx) - using python-pptx
- Legacy PowerPoint (.ppt) - converted to .pptx or returns empty

Usage:
    from src.document_parser import extract_pitch_deck_text

    text = extract_pitch_deck_text(Path("pitch_deck.pdf"))
"""

from __future__ import annotations

from pathlib import Path

from src.logging_config import get_logger

logger = get_logger(__name__)


def extract_text_from_pdf(file_path: Path) -> str:
    """Extract text content from a PDF file.

    Uses pdfplumber for better text layout handling compared to PyMuPDF.

    Args:
        file_path: Path to the PDF file.

    Returns:
        Extracted text content, or empty string if extraction fails.

    Example:
        >>> text = extract_text_from_pdf(Path("deck.pdf"))
        >>> "Investment Thesis" in text
        True
    """
    try:
        import pdfplumber

        text_parts: list[str] = []

        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        result = "\n\n".join(text_parts)
        logger.info(f"Extracted {len(result)} chars from PDF: {file_path.name}")
        return result

    except ImportError:
        logger.error("pdfplumber not installed. Run: uv add pdfplumber")
        return ""
    except Exception as e:
        logger.error(f"Failed to extract text from PDF {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text content from a PowerPoint file.

    Extracts text from all slides, including titles, body text,
    and text boxes.

    Args:
        file_path: Path to the PowerPoint file (.pptx).

    Returns:
        Extracted text content, or empty string if extraction fails.

    Example:
        >>> text = extract_text_from_pptx(Path("deck.pptx"))
        >>> "Market Opportunity" in text
        True
    """
    try:
        from pptx import Presentation

        text_parts: list[str] = []

        prs = Presentation(file_path)

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text.strip())

                # Also extract from tables if present
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text:
                                slide_texts.append(cell.text.strip())

            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        result = "\n\n".join(text_parts)
        logger.info(f"Extracted {len(result)} chars from PPTX: {file_path.name}")
        return result

    except ImportError:
        logger.error("python-pptx not installed. Run: uv add python-pptx")
        return ""
    except Exception as e:
        logger.error(f"Failed to extract text from PPTX {file_path}: {e}")
        return ""


def extract_pitch_deck_text(file_path: Path) -> str:
    """Extract text from a pitch deck file (PDF or PowerPoint).

    Dispatches to the appropriate extraction function based on file extension.

    Args:
        file_path: Path to the pitch deck file.

    Returns:
        Extracted text content, or empty string if:
        - File doesn't exist
        - Unsupported format
        - Extraction fails

    Supported extensions:
        - .pdf
        - .pptx
        - .ppt (legacy, limited support)

    Example:
        >>> text = extract_pitch_deck_text(Path("deck.pdf"))
        >>> len(text) > 0
        True
    """
    if not file_path.exists():
        logger.warning(f"File does not exist: {file_path}")
        return ""

    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    elif suffix == ".ppt":
        # Legacy .ppt format - python-pptx doesn't support it directly
        # In production, could convert via LibreOffice or similar
        logger.warning(f"Legacy .ppt format has limited support: {file_path}")
        # Try anyway - pptx might handle some files
        try:
            return extract_text_from_pptx(file_path)
        except Exception:
            logger.error(f"Cannot extract text from legacy .ppt file: {file_path}")
            return ""
    else:
        logger.warning(f"Unsupported file format: {suffix}")
        return ""


def get_supported_extensions() -> list[str]:
    """Get list of supported file extensions for pitch deck parsing.

    Returns:
        List of supported extensions including the dot prefix.
    """
    return [".pdf", ".pptx", ".ppt"]
